from pathlib import Path

from pptx import Presentation


def parse_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = shape.text.strip()
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_parts))
    return "\n\n".join(parts)
